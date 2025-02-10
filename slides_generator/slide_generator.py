from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from slides_generator.slide_data import SlideData
from slides_generator.font_size_calculation import FontSizeCalculation

class SlideGenerator:
    def __init__(self, presentation, slide_layout, slide_data : SlideData, path_logo : str = None):
        self.__presentation = presentation
        self.__slide_layout = slide_layout
        self.__slide_data = slide_data
        self.path_logo = path_logo

        self.calculate_font_size()

    def create_stanza_slide(self, insert_title : bool = False):
        self.__slide_pptx = self.__presentation.slides.add_slide(self.__slide_layout)
        self.set_background_color()

        if insert_title:
            self.add_title()
        self.add_content()

        if self.path_logo is not None:
            self.add_img_logo()

    def set_background_color(self):
        # Define o fundo preto do slide
        background = self.__slide_pptx.background
        fill = background.fill
        fill.solid()  # Configura para cor sólida
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Preto

    def add_title(self):
        if not hasattr(self, '__slide_pptx') or not hasattr(self, '__slide_data'):
            raise AttributeError("Atributos __slide_pptx ou __slide_data não encontrados")

        if not hasattr(self.__slide_data, 'title'):
            raise AttributeError("Atributo 'title' não encontrado em __slide_data")

        title_shape = self.__slide_pptx.shapes.title
        if title_shape is None:
            raise ValueError("O slide não possui um placeholder de título")
    
        title = self.__slide_pptx.shapes.title
        title.text = self.__slide_data.title

    def add_content(self):

        left = top = Inches(1)  # Ajuste as coordenadas conforme necessário
        width = Inches(8)  # Largura fixa do textbox
        height = Inches(2)  # Altura fixa do textbox

        txBox = self.__slide_pptx.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = self.__slide_data.stanza

        for p in tf.paragraphs:
            # Define a fonte e a cor (como no seu código anterior)
            p.font.size = self.__stanza_font_size
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Times New Roman'

         # Habilita a quebra de linha automática
        tf.word_wrap = True

    def add_img_logo(self):
        #img_path = 'img/logo.png'

        top = Inches(6.25)
        left = Inches(8.75)
        height = Inches(1)
        pic = self.__slide_pptx.shapes.add_picture(self.path_logo, left, top, height=height)

   
    def calculate_font_size(self):
        fsc = FontSizeCalculation()
        sfs = fsc.calcular_tamanho_fonte(self.__slide_data.stanza)
        self.__stanza_font_size = Pt(sfs)